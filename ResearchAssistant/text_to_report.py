from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from io import BytesIO
import datetime
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from openpyxl import Workbook
import os


def create_pdf(user_id, report_data, research_topic, delete_after_use=False):
    three_two_words = '_'.join(research_topic.split()[:3])  # Use underscores instead of spaces for the filename

    pdf_content_buffer = BytesIO()
    doc = SimpleDocTemplate(pdf_content_buffer, pagesize=letter)
    styles = getSampleStyleSheet()

    title_text = '<font size="16">{}</font>'.format(research_topic)
    title_paragraph = Paragraph(title_text, styles['Title'])
    elements = [title_paragraph]

    report_data_paragraph = Paragraph(report_data, styles['BodyText'])
    elements.append(report_data_paragraph)

    doc.build(elements)
    pdf_content = pdf_content_buffer.getvalue()
    pdf_content_buffer.close()

    # Create the "generate_report" directory if it doesn't exist
    if not os.path.exists("generate_report"):
        os.makedirs("generate_report")

    current_datetime = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    pdf_filename = f"generate_report/{three_two_words}_{user_id}_{current_datetime}.pdf"
    with open(pdf_filename, 'wb') as pdf_file:
        pdf_file.write(pdf_content)

    files = {'generated_report': (pdf_filename, open(pdf_filename, 'rb'), 'application/pdf')}

    if delete_after_use:
        os.remove(pdf_filename)
        del files['generated_report']
        return report_data, None  # Return None as the file is deleted

    return report_data, files


def create_docx(user_id, report_data, research_topic, delete_after_use=False):
    three_two_words = '_'.join(research_topic.split()[:3])

    doc = Document()
    doc.add_heading(research_topic, level=1)
    doc.add_paragraph(report_data)

    # Create the "generate_report" directory if it doesn't exist
    if not os.path.exists("generate_report"):
        os.makedirs("generate_report")

    current_datetime = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    docx_filename = f"generate_report/{three_two_words}_{user_id}_{current_datetime}.docx"
    doc.save(docx_filename)

    files = {'generated_report': (docx_filename, open(docx_filename, 'rb'),
                                  'application/vnd.openxmlformats-officedocument.wordprocessingml.document')}

    if delete_after_use:
        os.remove(docx_filename)
        del files['generated_report']
        return report_data, None  # Return None as the file is deleted

    return report_data, files


def create_pptx(user_id, report_data, research_topic, delete_after_use=False):
    three_two_words = '_'.join(research_topic.split()[:3])

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = research_topic

    content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(1.5), width=Inches(8), height=Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = report_data

    # Create the "generate_report" directory if it doesn't exist
    if not os.path.exists("generate_report"):
        os.makedirs("generate_report")

    current_datetime = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_filename = f"generate_report/{three_two_words}_{user_id}_{current_datetime}.pptx"
    prs.save(pptx_filename)

    files = {'generated_report': (pptx_filename, open(pptx_filename, 'rb'),
                                  'application/vnd.openxmlformats-officedocument.presentationml.presentation')}

    if delete_after_use:
        os.remove(pptx_filename)
        del files['generated_report']
        return report_data, None  # Return None as the file is deleted

    return report_data, files


def create_excel(user_id, report_data, research_topic, delete_after_use=False):
    three_two_words = '_'.join(research_topic.split()[:3])

    wb = Workbook()
    ws = wb.active
    ws.title = research_topic
    ws['A1'] = research_topic
    ws['A2'] = report_data

    # Create the "generate_report" directory if it doesn't exist
    if not os.path.exists("generate_report"):
        os.makedirs("generate_report")

    # Save the Excel file within the "generate_report" directory
    current_datetime = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    excel_filename = f"generate_report/{three_two_words}_{user_id}_{current_datetime}.xlsx"
    wb.save(excel_filename)

    files = {'generated_report': (excel_filename, open(excel_filename, 'rb'),
                                  'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')}

    if delete_after_use:
        os.remove(excel_filename)
        del files['generated_report']
        return report_data, None  # Return None as the file is deleted

    return report_data, files
